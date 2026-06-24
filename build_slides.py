#!/usr/bin/env python3
"""Results slideshow: minimal, white background, plain black Arial, plain titles,
explanatory confidence slides (including confidence per metal) and the comparisons."""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pandas as pd

BLACK = RGBColor(0, 0, 0); GRAY = RGBColor(0x40, 0x40, 0x40); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = 'Arial'; SW, SH = 13.333, 7.5
prs = Presentation(); prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

def slide():
    s = prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE; return s

def box(s, l, t, w, h):
    tf = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)).text_frame; tf.word_wrap = True
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'): setattr(tf, m, 0)
    return tf

def put(tf, text, size, color=BLACK, bold=False, italic=False, first=True, space_after=6):
    p = tf.paragraphs[0] if first else tf.add_paragraph(); p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(space_after)
    r = p.add_run(); r.text = text; f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic; f.color.rgb = color; f.name = FONT
    return p

def put2(tf, term, rest, size=16, first=True, space_after=13):
    p = tf.paragraphs[0] if first else tf.add_paragraph(); p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(space_after)
    for t, b in [(term, True), (rest, False)]:
        r = p.add_run(); r.text = t; r.font.size = Pt(size); r.font.bold = b; r.font.name = FONT; r.font.color.rgb = BLACK
    return p

def title(s, text):
    put(box(s, 0.7, 0.5, 12.0, 1.0), text, 29, BLACK, bold=True)

def add_img(s, path, top, box_w, box_h, center=True, left=0.9):
    w, h = Image.open(path).size; ar = w / h; W = box_w; H = W / ar
    if H > box_h: H = box_h; W = H * ar
    L = (SW - W) / 2 if center else left
    s.shapes.add_picture(path, Inches(L), Inches(top), Inches(W), Inches(H)); return top + H

def caption(s, text, top, w=11.9, l=0.7):
    put(box(s, l, top, w, 1.1), text, 13, GRAY, italic=True)

def fig_slide(fig, ttl, cap, box_h=4.5):
    s = slide(); title(s, ttl); bottom = add_img(s, f"figures/{fig}.png", 1.65, 11.6, box_h); caption(s, cap, bottom + 0.16); return s

# 1 title
s = slide(); tf = box(s, 0.7, 2.7, 12.0, 2.2)
put(tf, "Predicting logD for Rare-Earth Solvent Extraction", 38, BLACK, bold=True)
put(tf, "Using LLM-assisted coding to predict logD and judge how well an extractant separates two f-elements, with calibrated confidence and a comparison with Dr. Zhang's model", 17, BLACK, first=False)
put(box(s, 0.7, 5.1, 12.0, 0.5), "Skyler Epstein", 14, BLACK)

# 2 dataset
s = slide(); title(s, "Dataset")
put(box(s, 0.7, 1.5, 12.0, 1.7), "The goal is to predict how well an extractant separates two f-elements, with logD (the log of the distribution coefficient, D) as the target. The inputs are molecular descriptors, the SMILES string, molecular fingerprints, and environmental factors such as pH and temperature.", 16, BLACK)
x = 0.7
for num, lab in [("8,075", "measurements"), ("295", "distinct extractants"), ("28", "f-element metals"), ("~17.5", "log-unit range of logD")]:
    tf = box(s, x, 3.4, 3.0, 1.6); put(tf, num, 44, BLACK, bold=True); put(tf, lab, 14, GRAY, first=False); x += 3.05

# 3 two tasks
s = slide(); title(s, "Two prediction tasks")
put(box(s, 0.7, 1.5, 12.0, 0.8), "The tasks are kept apart so a molecule cannot leak across training and testing, which would inflate the score.", 16, BLACK)
L = box(s, 0.7, 2.6, 5.8, 3.4)
put(L, "Track A: screen a new molecule", 19, BLACK, bold=True); put(L, "molecule-grouped cross-validation", 13, GRAY, italic=True, first=False, space_after=12)
put(L, "R-squared 0.466     RMSE 1.148", 18, BLACK, bold=True, first=False, space_after=8); put(L, "Most confident 10 percent: R-squared 0.912, RMSE 0.493", 15, BLACK, first=False)
R = box(s, 7.0, 2.6, 5.8, 3.4)
put(R, "Track B: optimize conditions", 19, BLACK, bold=True); put(R, "random-row cross-validation", 13, GRAY, italic=True, first=False, space_after=12)
put(R, "R-squared 0.725     RMSE 0.823", 18, BLACK, bold=True, first=False, space_after=8); put(R, "Most confident 10 percent: R-squared 0.940, RMSE 0.341", 15, BLACK, first=False)

# 4 definitions
s = slide(); title(s, "Definitions")
tf = box(s, 0.7, 1.6, 12.0, 5.4)
put2(tf, "Confidence score:  ", "a second model estimates how far off each prediction is likely to be, so a small estimated error means high confidence.", first=True)
put2(tf, "Most confident 10 percent (or 25, or 50):  ", "the fraction of predictions with the smallest estimated error. Keeping only those gives much higher accuracy.", first=False)
put2(tf, "Uncertainty interval:  ", "instead of a single number, the model returns a range of logD values for each prediction.", first=False)
put2(tf, "90 percent interval:  ", "a range built to contain the true logD about 90 percent of the time. An 80 percent interval is narrower and contains it about 80 percent of the time.", first=False)
put2(tf, "Coverage:  ", "the fraction of points whose true value actually lands inside the interval. We measured 0.89 to 0.90 against the 90 percent target, so the ranges are honest.", first=False)

# 5-6 confidence curves
fig_slide("confidence_curve_a", "Accuracy versus confidence, Track A", "Moving right to left keeps fewer but more confident predictions. The most confident tenth (the tenth with the smallest estimated error) reaches R-squared 0.912 and RMSE 0.493, against 0.466 and 1.148 over all predictions.")
fig_slide("confidence_curve_b", "Accuracy versus confidence, Track B", "The same pattern holds for condition optimization: the most confident tenth reaches R-squared 0.940 and RMSE 0.341, against 0.725 and 0.823 over all predictions.")
# 7 pred vs actual
s = slide(); title(s, "Predicted versus actual logD")
add_img(s, "figures/pred_vs_actual_a.png", 1.75, 5.4, 4.7, center=False, left=1.2)
add_img(s, "figures/pred_vs_actual_b.png", 1.75, 5.4, 4.7, center=False, left=6.9)
caption(s, "Left: Track A, new molecules. Right: Track B, known molecules. Each point is a prediction; darker points are the more confident ones and lie closer to the diagonal, where predicted equals actual.", 6.65)

# 8-10 per metal
fig_slide("by_metal_accuracy", "Accuracy by metal", "Accuracy is not uniform. It is highest on Americium and Curium and lowest on Neptunium(V), Erbium, and Uranium(VI), so a single overall number hides large differences between metals.")
fig_slide("confidence_per_metal", "Confidence by metal", "Each bar is the model's median estimated error for a metal, which is its confidence signal, so a shorter bar means it is more sure. The model is most confident on Americium and Curium and least confident on Neptunium(V), Uranium(VI), and Erbium.", box_h=4.7)
fig_slide("by_metal_confidence_vs_error", "Confidence versus error, by metal", "The horizontal axis is the model's estimated error and the vertical axis is the actual error. They line up, so Americium and Curium sit low on both while Neptunium(V), Uranium(VI), and Erbium sit high on both, which is what lets the model flag the metals it cannot predict well.", box_h=4.3)

# 11-12 per pair
fig_slide("by_metal_pair_separation", "Separation between metal pairs", "A separation is the logD difference between two metals at the same conditions. Pairs far apart in logD such as Dysprosium and Ytterbium are predicted well, while tight adjacent pairs such as Erbium and Terbium are hardest, which is the case that matters most for industry.")
fig_slide("by_pair_confidence_vs_error", "Confidence by metal pair", "Confidence behaves the same way on pairs: the model's estimated error for a pair lines up with the actual separation error, so it is most confident on the easy far-apart pairs and least confident on hard adjacent pairs such as Erbium and Terbium.", box_h=4.3)

# 13 ensemble weights
fig_slide("ensemble_weights", "Ensemble weights", "Non-negative least squares keeps only the members that help, leaving ExtraTrees, LightGBM, and XGBoost on Track B and dropping the rest to a weight of zero.")
# 14 intervals
fig_slide("conformal_calibration", "Uncertainty intervals", "Each prediction comes with a range. A 90 percent interval is built to contain the true logD about 90 percent of the time, an 80 percent interval about 80 percent. The left bars show the measured coverage (0.89 to 0.90 and 0.79 to 0.80), so the ranges are honest; the right bars show the typical width, narrower on Track B.")

# 15 comparison setup
s = slide(); title(s, "Comparison with Dr. Zhang's model")
tf = box(s, 0.7, 1.8, 12.0, 4.0)
put(tf, "Both projects use the same 8,075-row f-element dataset (295 molecules, 28 metals), so the comparison is fair on the data.", 18, BLACK, space_after=14)
put(tf, "His supervised model is a 3-class XGBoost classifier that reports 72 percent accuracy on one 494-row held-out test, with no way to rank which predictions to trust.", 18, BLACK, first=False, space_after=14)
put(tf, "Ours predicts the continuous logD value with a confidence score and an uncertainty interval, and it can also be scored as a classifier on his exact task, which the next slides do.", 18, BLACK, first=False)

# 16 selective vs flat
fig_slide("zhang_comparison_selective", "Accuracy by confidence level, versus Zhang", "His 0.72 is a single value over all molecules. Ours can be ranked by confidence, so on the predictions it is most sure about it is far more accurate, which a single value cannot offer.", box_h=4.3)
# 17 his split
fig_slide("zhang_his_split_headtohead", "Our model on Dr. Zhang's test set", "On his exact 494-row test our model reaches 0.692, essentially matching his 0.72, and our confidence then sorts the reliable predictions to 0.863 over the top quarter and 1.000 over the most confident fifty.", box_h=4.3)

# 18 the 2x2 table
if os.path.exists("zhang_2x2_results.csv"):
    z = pd.read_csv("zhang_2x2_results.csv"); s = slide(); title(s, "Effect of model and split")
    tbl = s.shapes.add_table(3, 3, Inches(0.9), Inches(1.85), Inches(11.5), Inches(2.6)).table
    tbl.columns[0].width = Inches(4.3); tbl.columns[1].width = Inches(3.6); tbl.columns[2].width = Inches(3.6)
    grid = [["", "Our split (cross-validation)", "His split (494-row holdout)"],
            ["Our model (LightGBM)", f"{z.iloc[0]['our_split_CV']:.3f}", f"{z.iloc[0]['his_split_holdout']:.3f}"],
            ["His model (XGBoost)", f"{z.iloc[1]['our_split_CV']:.3f}", f"{z.iloc[1]['his_split_holdout']:.3f}  (his tuned: 0.72)"]]
    for i, row in enumerate(grid):
        for j, t in enumerate(row):
            cell = tbl.cell(i, j); cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            p = cell.text_frame.paragraphs[0]; r = p.add_run(); r.text = t
            r.font.size = Pt(15); r.font.bold = (i == 0 or j == 0); r.font.name = FONT; r.font.color.rgb = BLACK
    caption(s, "Same data and features throughout. On a common split the two models are about equal, ours slightly ahead, and the single held-out test reads a little higher than cross-validation for both, so the split explains most of the gap rather than the model. The bottom-right is an untuned reproduction of his XGBoost at 0.680, while his grid-tuned version reported 0.72. Our model on our own cross-validation is the top-left at 0.657.", 4.85)

# 19 his data is ours
fig_slide("our_data_vs_zhang_data", "Training on Dr. Zhang's data", "Training on his data reproduces our results, so his dataset is not better, it is the same; the edge is the feature engineering and the confidence layer, not the data.", box_h=4.1)

# 20 summary
s = slide(); title(s, "Summary")
tf = box(s, 0.7, 1.65, 12.0, 5.4)
pts = ["The dataset is shared with Dr. Zhang, so the difference is method, not data.",
       "On his exact test our model matches his accuracy, 0.69 against 0.72, and on a common split the two models are about equal.",
       "The confidence score estimates each prediction's error, so the most confident predictions are far more accurate, up to 1.000 on the most confident fifty of his test.",
       "Confidence shifts sensibly across metals and metal pairs, staying high where the model is accurate and dropping where it is not.",
       "Beyond a class, our model returns the continuous logD value and a calibrated uncertainty interval, neither of which his classifier provides.",
       "Next steps: a direct separation-factor model and replicate averaging to push past the current ceilings."]
for i, p in enumerate(pts): put(tf, p, 17, BLACK, first=(i == 0), space_after=12)

prs.save("REE_Results_Slides.pptx"); print("saved REE_Results_Slides.pptx with", len(prs.slides._sldIdLst), "slides")
