#!/usr/bin/env python3
"""Append the dataset-features slide and the free-energy (delta G) slides to the
existing deck, in the same minimal white / black Arial format, without altering any
of the existing slides. Opens the deck, adds slides at the end, and saves in place."""
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
DECK = "REE_logD_submission/docs/REE_Results_Slides_Final.pptx"
BLACK = RGBColor(0, 0, 0); GRAY = RGBColor(0x40, 0x40, 0x40); WHITE = RGBColor(0xFF, 0xFF, 0xFF); FONT = 'Arial'; SW, SH = 13.333, 7.5
prs = Presentation(DECK); BLANK = prs.slide_layouts[6]
def slide():
    s = prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE; return s
def box(s, l, t, w, h):
    tf = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)).text_frame; tf.word_wrap = True
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'): setattr(tf, m, 0)
    return tf
def put(tf, text, size, color=BLACK, bold=False, italic=False, first=True, space_after=6):
    p = tf.paragraphs[0] if first else tf.add_paragraph(); p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(space_after)
    r = p.add_run(); r.text = text; f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic; f.color.rgb = color; f.name = FONT; return p
def put2(tf, term, rest, size=16, first=True, space_after=13):
    p = tf.paragraphs[0] if first else tf.add_paragraph(); p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(space_after)
    for t, b in [(term, True), (rest, False)]:
        r = p.add_run(); r.text = t; r.font.size = Pt(size); r.font.bold = b; r.font.name = FONT; r.font.color.rgb = BLACK
    return p
def title(s, text): put(box(s, 0.7, 0.5, 12.0, 1.0), text, 29, BLACK, bold=True)
def add_img(s, path, top, box_w, box_h):
    w, h = Image.open(path).size; ar = w / h; W = box_w; H = W / ar
    if H > box_h: H = box_h; W = H * ar
    s.shapes.add_picture(path, Inches((SW - W) / 2), Inches(top), Inches(W), Inches(H)); return top + H
def caption(s, text, top): put(box(s, 0.7, top, 11.9, 1.1), text, 13, GRAY, italic=True)
n0 = len(prs.slides._sldIdLst)
# dataset features
s = slide(); title(s, "Dataset features and why")
tf = box(s, 0.7, 1.6, 12.0, 5.4)
put2(tf, "Extractant structure:  ", "SMILES, a Morgan fingerprint, RDKit descriptors, and donor-atom counts (O, N, S, P) with logP. These set how, and how strongly, the extractant binds a metal.", size=15, first=True, space_after=14)
put2(tf, "Metal descriptors:  ", "ionic radius, charge and oxidation state, ionization energies, electronegativity, and atomic number. These carry the lanthanide-contraction trend that makes neighboring rare earths hard to separate.", size=15, first=False, space_after=14)
put2(tf, "Acid and conditions:  ", "acid type and concentration, temperature, and extractant concentration. logD is an equilibrium, so it shifts with these as much as with the molecule.", size=15, first=False, space_after=14)
put2(tf, "Diluent (solvent):  ", "molar mass, logP, boiling and melting points, density, solubility, and dipole moment. The diluent sets the organic-phase environment around the extractant.", size=15, first=False, space_after=14)
# free energy framing
s = slide(); title(s, "Predicting the free energy of extraction")
tf = box(s, 0.7, 1.6, 12.0, 5.4)
put2(tf, "Target:  ", "delta G = -2.303 R T logD, in kJ/mol, so a favorable extraction is a negative, spontaneous free energy. Predicted from the extractant structure and the metal, with the reaction conditions left out.", size=15, first=True, space_after=16)
put2(tf, "Why drop the conditions:  ", "the free energy is meant to be a property of the extractant and the metal, not of the acid concentration or temperature, so the model predicts the thermodynamics from structure alone.", size=15, first=False, space_after=16)
put2(tf, "Two framings:  ", "per-row keeps every measurement; per-pair averages delta G to one value per extractant, metal, acid, and diluent system and is condition-independent, which is the cleaner target.", size=15, first=False, space_after=16)
# free energy results figure
s = slide(); title(s, "Free energy (delta G) results")
bottom = add_img(s, "figures/dg.png", 1.7, 11.6, 4.3)
caption(s, "Per-pair is the better target (R-squared about 0.46, RMSE 6.3 kJ/mol) than per-row (0.28), molecule-grouped so the extractants are new. The 0.473 shown elsewhere was a selected maximum; the mean over shuffled splits is 0.461 plus or minus 0.009. Confidence is strong: the most confident tenth drop RMSE to 3.6 kJ/mol.", bottom + 0.18)
# delta G ensemble sweep
s = slide(); title(s, "Delta G ensemble sweep")
tf = box(s, 0.7, 1.6, 12.0, 5.4)
put2(tf, "Best model:  ", "RandomForest, R-squared about 0.46 plus or minus 0.01 (RMSE 6.3 kJ/mol) for the per-pair free energy; 0.473 was the maximum of the search, the mean over splits is 0.461. The bagged-tree models beat the gradient boosters on this small, wide table of 2,273 systems.", size=15, first=True, space_after=16)
put2(tf, "The field:  ", "XGBoost 0.444, CatBoost 0.441, HistGB 0.426, LightGBM 0.424, ExtraTrees 0.419, and a linear Ridge near zero, so the signal is nonlinear.", size=15, first=False, space_after=16)
put2(tf, "Stacking:  ", "an NNLS stack, scored with its own cross-validation, reaches 0.474, a tie with RandomForest alone, so stacking is not used and the free-energy model is a plain RandomForest.", size=15, first=False, space_after=16)
# UCB and active-learning graphs
s = slide(); title(s, "UCB and active-learning results")
bottom = add_img(s, "figures/ucb_graphs.png", 1.75, 12.4, 3.9)
caption(s, "Left: sampling the uncertain points builds the best model. Middle: ranking by prediction (greedy) or UCB finds the strongest extractants fastest. Right: for a one-shot pick on delta G, greedy reaches the most negative, strongest extractants, UCB sits in the middle, and random gets nothing. UCB does not beat the simpler choices here.", bottom + 0.18)
# active analysis (calibrated intervals)
s = slide(); title(s, "Active analysis: calibrated intervals")
bottom = add_img(s, "figures/active_analysis.png", 1.75, 12.4, 3.9)
caption(s, "Left: the delta G intervals are calibrated, target coverage matches empirical coverage at every level. Middle: keeping only the most confident predictions raises R-squared from 0.47 to 0.78 and cuts RMSE from 6.3 to 3.6 kJ/mol. Right: the most confident quarter (blue) hug the diagonal. Calibrated, heterogeneous intervals are what let the model say which predictions to trust and which to test.", bottom + 0.18)
# honest numbers after the audit
s = slide(); title(s, "Honest numbers after the audit")
tf = box(s, 0.7, 1.5, 12.0, 5.6)
put2(tf, "Track A (new molecule):  ", "logD R-squared 0.466, unchanged and honest.", size=15, first=True, space_after=13)
put2(tf, "Track B (known molecule):  ", "honest R-squared about 0.61, interpolating conditions for a molecule already seen. The 0.725 figure was a random-row split that memorizes replicates; on genuinely new molecules these features fall to about 0.44.", size=15, first=False, space_after=13)
put2(tf, "Free energy (delta G):  ", "R-squared about 0.46 plus or minus 0.01; the 0.473 was a selected maximum of a model search.", size=15, first=False, space_after=13)
put2(tf, "Confidence:  ", "the 90 percent intervals are calibrated (they cover 90 percent); the most-confident-slice figures are operating points, not overall accuracy.", size=15, first=False, space_after=13)
put2(tf, "Versus Dr. Zhang:  ", "about equal on the same data and split (0.657 vs 0.648); the split, not the model, drives his 0.72 headline.", size=15, first=False, space_after=13)
# active-analysis trends
s = slide(); title(s, "Active analysis: good picks or mid?")
bottom = add_img(s, "figures/active_analysis_trends.png", 1.55, 11.8, 4.6)
caption(s, "Greedy selection puts 81 percent of its top picks in the strongest tertile (vs 37 percent random) and beats each metal's average, but it captures only about 62 percent of the achievable strength and about 57 percent of its top-decile picks miss the true elite. Its confident picks are the trustworthy ones.", bottom + 0.16)
# chosen-extractant trends
s = slide(); title(s, "What the chosen extractants share")
bottom = add_img(s, "figures/picks_trends.png", 1.9, 12.2, 4.0)
caption(s, "The picks concentrate on americium and the trivalent lanthanides and actinides (americium 2.85x over-represented, 30 percent of picks), and lean toward less-aromatic, more aliphatic extractants, consistent with classic organophosphorus f-element extractants. Effect sizes are modest.", bottom + 0.18)
prs.save(DECK)
print(f"appended {len(prs.slides._sldIdLst) - n0} slides; deck now {len(prs.slides._sldIdLst)} slides")
