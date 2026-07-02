#!/usr/bin/env python3
"""Append the AutoData / selectivity / new-extractant-confidence slides to the CURRENT deck,
append-only, preserving every existing slide (never rebuild)."""
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
DECK = "docs/REE_Results_Slides_Final.pptx"
BLACK = RGBColor(0, 0, 0); GRAY = RGBColor(0x40, 0x40, 0x40); WHITE = RGBColor(0xFF, 0xFF, 0xFF); FONT = 'Arial'; SW, SH = 13.333, 7.5
prs = Presentation(DECK); BLANK = prs.slide_layouts[6]
def slide():
    s = prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE; return s
def box(s, l, t, w, h):
    tf = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)).text_frame; tf.word_wrap = True
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'): setattr(tf, m, 0)
    return tf
def put(tf, text, size, color=BLACK, bold=False, italic=False):
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT; r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic; f.color.rgb = color; f.name = FONT
def title(s, text): put(box(s, 0.7, 0.5, 12.0, 1.0), text, 29, BLACK, bold=True)
def add_img(s, path, top, box_w, box_h):
    w, h = Image.open(path).size; ar = w / h; W = box_w; H = W / ar
    if H > box_h: H = box_h; W = H * ar
    s.shapes.add_picture(path, Inches((SW - W) / 2), Inches(top), Inches(W), Inches(H)); return top + H
def caption(s, text, top): put(box(s, 0.7, top, 11.9, 1.2), text, 13, GRAY, italic=True)
n0 = len(prs.slides._sldIdLst)

s = slide(); title(s, "AutoData: generating candidate extractants")
caption(s, "An agentic generate-score-refine loop proposes novel extractant molecules. It rediscovers the right soft-donor families (dithiophosphinic acids, thiopicolinamides) but does not beat the best known extractants: novel predicted separation p90 near a factor of 10, the in-sample known near 22. A candidate generator for the wet lab, not a labeler.", add_img(s, "figures/autodata_pool.png", 1.7, 12.3, 4.2) + 0.14)

s = slide(); title(s, "Does it separate two metals, or just extract everything?")
caption(s, "The honest check: separation needs selectivity, not general strength. On known extractants the model captures real selectivity (uncorrelated with strength, correct soft-donor mechanism). On novel molecules it partly collapses into strength, so the ranking is strength-adjusted.", add_img(s, "figures/autodata_selectivity.png", 1.7, 12.3, 4.2) + 0.14)

s = slide(); title(s, "Separation as the target: two confidence-gated shortlists")
caption(s, "With separation set as the search target, each molecule is scored by its best f-element pair. EXPLORE keeps pairs with at least 5 measured extractants (higher factors, thin data); TRUSTED keeps at least 20 (real support). Both are gated to the most confident top 25 and top 10 percent.", add_img(s, "figures/autodata_bestpair_confident.png", 1.7, 12.3, 4.2) + 0.14)

s = slide(); title(s, "A confidence algorithm that works for new extractants")
caption(s, "Bake-off winner: bootstrap-bagged ensemble disagreement. It is the only method that sharpens new-extractant separation (top-10 percent direction 0.66 to 0.81, signed R-squared about 0.00 to 0.28, calibration 0.20 to 0.31). The applicability-domain distances expected to win lost. The known-extractant confidence layer is unchanged.", add_img(s, "figures/newext_confidence_bakeoff.png", 1.7, 12.3, 4.2) + 0.14)

prs.save(DECK)
print(f"appended {len(prs.slides._sldIdLst) - n0} slides; deck now {len(prs.slides._sldIdLst)}")
