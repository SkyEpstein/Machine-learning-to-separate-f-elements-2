#!/usr/bin/env python3
"""Append the three separation slides (including the confidence-forward Zhang
comparison) to the CURRENT deck, without rebuilding, so the other session's edits
(e.g. the recomputed 0.874) are preserved."""
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
DECK = "REE_logD_submission/docs/REE_Results_Slides_Final.pptx"
BLACK = RGBColor(0, 0, 0); GRAY = RGBColor(0x40, 0x40, 0x40); WHITE = RGBColor(0xFF, 0xFF, 0xFF); FONT = 'Arial'; SW, SH = 13.333, 7.5
prs = Presentation(DECK); BLANK = prs.slide_layouts[6]
def slide():
    s = prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE; return s
def box(s, l, t, w, h):
    tf = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)).text_frame; tf.word_wrap = True
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'): setattr(tf, m, 0)
    return tf
def put(tf, text, size, color=BLACK, bold=False, italic=False):
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT; r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic; f.color.rgb = color; f.name = FONT
def title(s, text): put(box(s, 0.7, 0.5, 12.0, 1.0), text, 29, BLACK, bold=True)
def add_img(s, path, top, box_w, box_h):
    w, h = Image.open(path).size; ar = w / h; W = box_w; H = W / ar
    if H > box_h: H = box_h; W = H * ar
    s.shapes.add_picture(path, Inches((SW - W) / 2), Inches(top), Inches(W), Inches(H)); return top + H
def caption(s, text, top): put(box(s, 0.7, top, 11.9, 1.2), text, 13, GRAY, italic=True)
n0 = len(prs.slides._sldIdLst)
s = slide(); title(s, "Separation between two f-elements")
caption(s, "Separation is obtained by differencing the logD model. The order, which f-element extracts more, is predicted moderately (direction 0.73 known, 0.66 new), even for near-identical neighbors; the size of the gap is not (magnitude R-squared 0.17 known, negative for new).", add_img(s, "figures/sep_factor_eval.png", 1.7, 12.3, 4.2) + 0.16)
s = slide(); title(s, "Separation: confidence helps known extractants")
caption(s, "Keeping only the most confident pairs improves a known extractant (direction 0.73 to 0.79, error dropping below the no-skill line) but not a new one (direction slips, error stays flat). The confidence layer earns its keep on known extractants.", add_img(s, "figures/sep_factor_confidence.png", 1.7, 12.3, 4.2) + 0.16)
s = slide(); title(s, "Separation vs Dr. Zhang, with our confidence")
caption(s, "The confidence layer is our contribution and is model-agnostic. His published model reports one flat number; applied to either base model our model leads on the confident slice (signed R-squared 0.59 vs 0.47, direction 0.79 vs 0.70, known extractant).", add_img(s, "figures/sep_zhang_confidence.png", 1.7, 12.3, 4.2) + 0.16)
prs.save(DECK)
print(f"appended {len(prs.slides._sldIdLst) - n0} slides; deck now {len(prs.slides._sldIdLst)}")
