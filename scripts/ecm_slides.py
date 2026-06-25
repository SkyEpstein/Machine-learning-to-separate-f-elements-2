#!/usr/bin/env python3
"""Append the dataset-features slide and the ECM (free energy) slides to the
existing deck, in the same minimal white / black Arial format, without altering any
of the existing slides. Opens the deck, adds slides at the end, and saves in place."""
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
DECK = "REE_logD_submission/docs/REE_Results_Slides_Final.pptx"
BLACK = RGBColor(0, 0, 0); GRAY = RGBColor(0x40, 0x40, 0x40); WHITE = RGBColor(0xFF, 0xFF, 0xFF); FONT = 'Arial'; SW, SH = 13.333, 7.5
prs = Presentation(DECK); BLANK = prs.slide_layouts[6]
def slide():
    s = prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE; return s
def box(s, l, t, w, h):
    tf = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)).text_frame; tf.word_wrap = True
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'): setattr(tf, m, 0)
    return tf
def put(tf, text, size, color=BLACK, bold=False, italic=False, first=True, space_after=6):
    p = tf.paragraphs[0] if first else tf.add_paragraph(); p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(space_after)
    r = p.add_run(); r.text = text; f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic; f.color.rgb = color; f.name = FONT; return p
def put2(tf, term, rest, size=16, first=True, space_after=13):
    p = tf.paragraphs[0] if first else tf.add_paragraph(); p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(space_after)
    for t, b in [(term, True), (rest, False)]:
        r = p.add_run(); r.text = t; r.font.size = Pt(size); r.font.bold = b; r.font.name = FONT; r.font.color.rgb = BLACK
    return p
def title(s, text): put(box(s, 0.7, 0.5, 12.0, 1.0), text, 29, BLACK, bold=True)
def add_img(s, path, top, box_w, box_h):
    w, h = Image.open(path).size; ar = w / h; W = box_w; H = W / ar
    if H > box_h: H = box_h; W = H * ar
    s.shapes.add_picture(path, Inches((SW - W) / 2), Inches(top), Inches(W), Inches(H)); return top + H
def caption(s, text, top): put(box(s, 0.7, top, 11.9, 1.1), text, 13, GRAY, italic=True)
n0 = len(prs.slides._sldIdLst)
# dataset features
s = slide(); title(s, "Dataset features and why")
tf = box(s, 0.7, 1.6, 12.0, 5.4)
put2(tf, "Extractant structure:  ", "SMILES, a Morgan fingerprint, RDKit descriptors, and donor-atom counts (O, N, S, P) with logP. These set how, and how strongly, the extractant binds a metal.", size=15, first=True, space_after=14)
put2(tf, "Metal descriptors:  ", "ionic radius, charge and oxidation state, ionization energies, electronegativity, and atomic number. These carry the lanthanide-contraction trend that makes neighboring rare earths hard to separate.", size=15, first=False, space_after=14)
put2(tf, "Acid and conditions:  ", "acid type and concentration, temperature, and extractant concentration. logD is an equilibrium, so it shifts with these as much as with the molecule.", size=15, first=False, space_after=14)
put2(tf, "Diluent (solvent):  ", "molar mass, logP, boiling and melting points, density, solubility, and dipole moment. The diluent sets the organic-phase environment around the extractant.", size=15, first=False, space_after=14)
# ECM framing
s = slide(); title(s, "ECM: predicting the free energy of extraction")
tf = box(s, 0.7, 1.6, 12.0, 5.4)
put2(tf, "Target:  ", "delta G = -2.303 R T logD, in kJ/mol, so a favorable extraction is a negative, spontaneous free energy. Predicted from the extractant structure and the metal, with the reaction conditions left out.", size=15, first=True, space_after=16)
put2(tf, "Why drop the conditions:  ", "the free energy is meant to be a property of the extractant and the metal, not of the acid concentration or temperature, so the model predicts the thermodynamics from structure alone.", size=15, first=False, space_after=16)
put2(tf, "Two framings:  ", "per-row keeps every measurement; per-pair averages delta G to one value per extractant, metal, acid, and diluent system and is condition-independent, which is the cleaner target.", size=15, first=False, space_after=16)
# ECM results figure
s = slide(); title(s, "ECM results")
bottom = add_img(s, "figures/ecm.png", 1.7, 11.6, 4.3)
caption(s, "Per-pair is the better target (R-squared 0.42, RMSE 6.6 kJ/mol) than per-row (0.27), evaluated with molecule-grouped cross-validation so the extractants are new. Confidence is strong: the most confident quarter reach R-squared 0.68 and the top tenth 0.81 (RMSE 3.2 kJ/mol).", bottom + 0.18)
prs.save(DECK)
print(f"appended {len(prs.slides._sldIdLst) - n0} slides; deck now {len(prs.slides._sldIdLst)} slides")
